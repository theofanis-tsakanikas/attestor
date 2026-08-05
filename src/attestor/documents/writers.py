"""Emitting the three artefacts.

Word for the statement, Excel for the datapoint annex, PowerPoint for the board. The
regulation asks for the first; the auditor works from the second; the third is how the
decision actually gets made.

Every writer takes the *same* `RenderedDocument` and writes the manifest beside the file.
The provenance gate then reads the file back — not the object that produced it — because a
manifest is a claim about a document, and a claim nobody checks against the artefact is a
comment.

Determinism note: Office formats embed a creation timestamp and a random document id, so two
runs produce different bytes for identical content. Rather than fight the format, claim 4 is
asserted at the level that matters — identical values and identical lineage hashes — and the
byte-level check applies to the manifest JSON, which *is* stable.
"""

from __future__ import annotations

import contextlib
import json
from pathlib import Path

from attestor.documents.render import RenderedDocument

_HEADING_SIZES = {1: 16, 2: 13, 3: 11, 4: 10}


def write_manifest(document: RenderedDocument, path: Path) -> Path:
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(
        json.dumps(document.manifest.as_dict(), indent=2, sort_keys=True, ensure_ascii=False),
        encoding="utf-8",
    )
    return path


def write(document: RenderedDocument, path: Path) -> Path:
    """Dispatch on the template's declared artefact."""
    match document.template.artefact:
        case "docx":
            return write_docx(document, path)
        case "xlsx":
            return write_xlsx(document, path)
        case "pptx":
            return write_pptx(document, path)
    raise ValueError(f"unknown artefact {document.template.artefact!r}")  # pragma: no cover


# ── Word ─────────────────────────────────────────────────────────────────────


def write_docx(document: RenderedDocument, path: Path) -> Path:
    from docx import Document  # noqa: PLC0415 — heavy import, only when writing
    from docx.shared import Pt  # noqa: PLC0415

    doc = Document()
    for section in document.sections:
        for block in section.blocks:
            match block.kind:
                case "title":
                    doc.add_heading(block.text, level=0)
                case "subtitle":
                    paragraph = doc.add_paragraph()
                    run = paragraph.add_run(block.text)
                    run.italic = True
                case "heading":
                    doc.add_heading(block.text, level=min(block.level, 4))
                case "paragraph" | "note":
                    paragraph = doc.add_paragraph()
                    for run in block.runs:
                        added = paragraph.add_run(run.text)
                        if block.kind == "note":
                            added.italic = True
                            added.font.size = Pt(9)
                case "bullets":
                    for row in block.rows:
                        doc.add_paragraph(row[0], style="List Bullet")
                case "table":
                    _docx_table(doc, block)
                case "page_break":
                    doc.add_page_break()
    path.parent.mkdir(parents=True, exist_ok=True)
    doc.save(str(path))
    return path


def _docx_table(doc, block) -> None:
    if not block.rows:
        return
    table = doc.add_table(rows=1, cols=len(block.headers or block.rows[0]))
    table.style = "Light Grid Accent 1"
    for index, header in enumerate(block.headers):
        table.rows[0].cells[index].text = header
    for row in block.rows:
        cells = table.add_row().cells
        for index, value in enumerate(row):
            cells[index].text = str(value)


# ── Excel ────────────────────────────────────────────────────────────────────


def write_xlsx(document: RenderedDocument, path: Path) -> Path:
    """The auditor annex.

    Values are written as text, not as numbers. That looks wrong for a spreadsheet and is
    deliberate: a float cell invites a reader to recompute a disclosure in Excel, and the
    figure that was signed for is the one the contract rounded. The trace sheet is what the
    auditor actually walks.
    """
    from openpyxl import Workbook  # noqa: PLC0415
    from openpyxl.styles import Alignment, Font  # noqa: PLC0415

    workbook = Workbook()
    workbook.remove(workbook.active)

    for section in document.sections:
        tables = [b for b in section.blocks if b.kind == "table" and b.rows]
        if not tables:
            continue
        sheet = workbook.create_sheet(_sheet_name(section.title or document.template.title))
        cursor = 1
        for block in tables:
            for column, header in enumerate(block.headers, start=1):
                cell = sheet.cell(row=cursor, column=column, value=header)
                cell.font = Font(bold=True)
            cursor += 1
            for row in block.rows:
                for column, value in enumerate(row, start=1):
                    cell = sheet.cell(row=cursor, column=column, value=str(value))
                    cell.alignment = Alignment(vertical="top", wrap_text=True)
                cursor += 1
            cursor += 1
        for column in "ABCDE":
            sheet.column_dimensions[column].width = 34

    if not workbook.sheetnames:  # pragma: no cover — a table-less xlsx template
        workbook.create_sheet("empty")
    path.parent.mkdir(parents=True, exist_ok=True)
    workbook.save(str(path))
    return path


def _sheet_name(title: str) -> str:
    cleaned = "".join(c for c in title if c not in "[]:*?/\\")
    return (cleaned or "sheet")[:31]


# ── PowerPoint ───────────────────────────────────────────────────────────────


def write_pptx(document: RenderedDocument, path: Path) -> Path:
    from pptx import Presentation  # noqa: PLC0415
    from pptx.util import Inches, Pt  # noqa: PLC0415

    presentation = Presentation()
    blank = presentation.slide_layouts[6]
    title_layout = presentation.slide_layouts[0]

    first = document.sections[0] if document.sections else None
    if first and first.blocks:
        slide = presentation.slides.add_slide(title_layout)
        slide.shapes.title.text = first.blocks[0].text
        if len(first.blocks) > 1 and slide.placeholders:
            # A layout without a subtitle placeholder is a cosmetic loss, not a data one.
            with contextlib.suppress(KeyError, IndexError):
                slide.placeholders[1].text = first.blocks[1].text

    for section in document.sections[1:]:
        slide = presentation.slides.add_slide(blank)
        box = slide.shapes.add_textbox(Inches(0.6), Inches(0.5), Inches(8.8), Inches(0.9))
        heading = box.text_frame.paragraphs[0]
        heading.text = section.title
        heading.runs[0].font.size = Pt(28)
        heading.runs[0].font.bold = True

        body = slide.shapes.add_textbox(Inches(0.6), Inches(1.6), Inches(8.8), Inches(5.0))
        frame = body.text_frame
        frame.word_wrap = True
        first_line = True
        for block in section.blocks:
            if block.kind == "heading":
                continue
            lines = (
                [row[0] for row in block.rows]
                if block.kind == "bullets"
                else [" · ".join(str(c) for c in row) for row in block.rows]
                if block.kind == "table"
                else [block.text]
            )
            for line in lines:
                if not line:
                    continue
                paragraph = frame.paragraphs[0] if first_line else frame.add_paragraph()
                paragraph.text = line
                paragraph.font.size = Pt(14)
                first_line = False

    path.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(str(path))
    return path
